import openai
import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
from dotenv import load_dotenv
from openai import OpenAI

load_dotenv()

client = OpenAI(
    api_key=os.environ["OPENAI_API_KEY"],
)

# custom formating
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)


def generate_slide_title(topic):
    prompt = f"Generate 5 slide titles for the given topic '{topic}."
    response = client.completions.create(
        model="gpt-3.5-turbo-0613", prompt=prompt, stream=200
    )
    return response["choices"][0]["text"].split("\n")


def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}."
    response = client.completions.create(
        model="gpt-3.5-turbo-0613", prompt=prompt, stream=500
    )
    return response["choices"][0]["text"]


def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layout[1]
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        slide.shapes.title.text_frame.pragaraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    prs.save(f"generated_ppt/{topic}_presentaion.pptx")


def main():
    st.title("Text to PPT Generation using LLM")
    topic = st.text_input("ENter the topic you want to Generate the PPT presentation")
    generate_button = st.button("Generate PPT")

    if generate_button and topic:
        st.info("Generating PPT presentation.....Please wait.")
        slide_titles = generate_slide_title(topic)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != ""]
        print("Slide titles: ", filtered_slide_titles)
        slide_contents = [
            generate_slide_content(topic) for title in filtered_slide_titles
        ]
        print("Slide contents: ", slide_contents)
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation Generated Successfully.")

        st.success("Presantation Generated Successfully")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)


def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentaion.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()

    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentation.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the Powerpoint presentation </a>'


if __name__ == "__main__":
    main()
